

#=================================================
# backend/app/services/ppt_service.py
# backend/app/services/ppt_service.py
# backend/app/services/ppt_service.py
# import os
# import json
# import uuid
# from datetime import datetime
# from typing import Any, Dict, List, Optional

# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.shapes import MSO_SHAPE
# from pptx.dml.color import RGBColor

# from reportlab.lib.pagesizes import landscape, A4
# from reportlab.pdfgen import canvas

# from app import config
# from io import BytesIO
# from app.utils import coerce_json
# from app.services.export_service import save_export_record

# OUT_DIR = config.OUTPUT_PPT_DIR
# os.makedirs(OUT_DIR, exist_ok=True)


# def _add_bullets_to_slide(slide, points: List[str]) -> None:
#     """
#     Add bullet points to a slide using a textbox. Ensures the first paragraph
#     is not left blank (TextFrame has one empty paragraph by default).
#     """
#     if not points:
#         return

#     left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
#     textbox = slide.shapes.add_textbox(left, top, width, height)
#     tf = textbox.text_frame
#     tf.clear()  # leaves a single empty paragraph

#     # fill first paragraph, then add remaining as new paragraphs
#     first_para = tf.paragraphs[0]
#     first_para.text = points[0]
#     first_para.font.size = Pt(18)
#     first_para.level = 0

#     for p in points[1:]:
#         para = tf.add_paragraph()
#         para.text = p
#         para.font.size = Pt(18)
#         para.level = 0


# def _add_optional_image(slide, image_path: Optional[str]) -> None:
#     """
#     Add an image to the right side of the slide if the file exists,
#     otherwise add a gray placeholder rectangle.
#     """
#     if not image_path:
#         return

#     try:
#         if os.path.exists(image_path):
#             slide.shapes.add_picture(image_path, Inches(6), Inches(1.5), width=Inches(3))
#         else:
#             # Placeholder if the path does not exist
#             shape = slide.shapes.add_shape(
#                 MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.5), Inches(3), Inches(3)
#             )
#             shape.fill.solid()
#             shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
#             shape.text = "Image"
#     except Exception:
#         # Fallback placeholder if add_picture fails
#         shape = slide.shapes.add_shape(
#             MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.5), Inches(3), Inches(3)
#         )
#         shape.fill.solid()
#         shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
#         shape.text = "Image"


# def create_presentation(
#     slides_json: Any,
#     template: Optional[str] = None,
#     num_slides: Optional[int] = None,
#     output_format: str = "pptx",
#     ttl_minutes: int = 30,
#     topic: Optional[str] = None,
#     audience: Optional[str] = None,
# ) -> Dict[str, Any]:
#     """
#     Build a PPTX or PDF from structured slide JSON, record the export in DB, and return meta.

#     slides_json: dict or JSON str with shape:
#       {
#         "slides": [
#           {
#             "title": "Slide Title",
#             "content": ["point 1", "point 2"],
#             "speaker_notes": "optional notes",
#             "_image_path": "/abs/or/rel/path.png"  # optional
#           },
#           ...
#         ]
#       }

#     Returns:
#       { "file_id": <uuid str>, "path": <absolute_path>, "download_url": <url>, "record_id": <db id> }
#     """
#     # Accept dict or str
#     if isinstance(slides_json, str):
#         slides_obj = coerce_json(slides_json)
#     else:
#         slides_obj = slides_json or {}

#     slides = slides_obj.get("slides", [])
#     if num_slides:
#         slides = slides[:num_slides]

#     # fallback metadata from slides_json if not provided explicitly
#     meta_topic = topic or slides_obj.get("topic") or "Unknown"
#     meta_audience = audience or slides_obj.get("audience") or "General Public"

#     # Load template if provided
#     if template and os.path.exists(template):
#         prs = Presentation(template)
#     else:
#         prs = Presentation()

#     # Build PPT slides
#     for s in slides:
#         title: str = s.get("title", "")
#         content: List[str] = s.get("content", []) or []
#         notes: str = s.get("speaker_notes", "")
#         image_path: Optional[str] = s.get("_image_path") or s.get("image_path") or None

#         layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
#         slide = prs.slides.add_slide(layout)

#         # Title
#         if slide.shapes.title:
#             slide.shapes.title.text = title

#         # Bullets
#         _add_bullets_to_slide(slide, content)

#         # Image
#         _add_optional_image(slide, image_path)

#         # Notes
#         try:
#             ns = slide.notes_slide
#             ns.notes_text_frame.text = notes or ""
#         except Exception:
#             # notes_slide may not exist on some layouts if notes master is unavailable
#             pass

#     # Generate unique file id & names
#     file_uuid = str(uuid.uuid4())
#     timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
#     base = f"presentation_{timestamp}_{file_uuid}"

#     # Save PPTX or PDF
#     if output_format.lower() == "pptx":
#         pptx_path = os.path.join(OUT_DIR, base + ".pptx")
#         prs.save(pptx_path)

#         # record in DB via export_service (this calls app.db.save_export_record)
#         rec = save_and_record_export(pptx_path, topic=meta_topic, audience=meta_audience, output_format="pptx")
#         # rec: {"record_id": int, "download_url": str, "file_path": str}
#         return {
#             "file_id": file_uuid,
#             "path": pptx_path,
#             "download_url": rec["download_url"],
#             "record_id": rec["record_id"]
#         }

#     elif output_format.lower() == "pdf":
#         pdf_path = os.path.join(OUT_DIR, base + ".pdf")
#         width, height = landscape(A4)
#         c = canvas.Canvas(pdf_path, pagesize=landscape(A4))

#         for s in slides:
#             title = s.get("title", "")
#             content = s.get("content", []) or []

#             # Title
#             c.setFont("Helvetica-Bold", 24)
#             c.drawString(40, height - 80, title)

#             # Bullets
#             c.setFont("Helvetica", 14)
#             y = height - 120
#             for p in content:
#                 c.drawString(60, y, u"\u2022 " + p)
#                 y -= 24
#                 if y < 60:
#                     # new simple page if overflow (keep it simple & consistent with prior logic)
#                     c.showPage()
#                     c.setFont("Helvetica", 14)
#                     y = height - 120
#             c.showPage()

#         c.save()
#         rec = save_and_record_export(pdf_path, topic=meta_topic, audience=meta_audience, output_format="pdf")
#         return {
#             "file_id": file_uuid,
#             "path": pdf_path,
#             "download_url": rec["download_url"],
#             "record_id": rec["record_id"]
#         }

#     else:
#         raise ValueError(f"Unsupported format: {output_format}")

#-=======================================================================
import os
import uuid
from datetime import datetime
from typing import Any, Dict, List, Optional
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import landscape, A4
from reportlab.pdfgen import canvas

from app import config
from app.services.export_service import save_and_record_export
from app.utils import coerce_json

OUT_DIR = config.OUTPUT_PPT_DIR
os.makedirs(OUT_DIR, exist_ok=True)

def _add_bullets_to_slide(slide, points: List[str]) -> None:
    if not points:
        return
    left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    first_para = tf.paragraphs[0]
    first_para.text = points[0]
    first_para.font.size = Pt(18)
    first_para.level = 0
    for p in points[1:]:
        para = tf.add_paragraph()
        para.text = p
        para.font.size = Pt(18)
        para.level = 0

def _add_optional_image(slide, image_path: Optional[str]) -> None:
    if not image_path:
        return
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(6), Inches(1.5), width=Inches(3))
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.5), Inches(3), Inches(3))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
            shape.text = "Image"
    except Exception:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.5), Inches(3), Inches(3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
        shape.text = "Image"

def create_presentation(
    slides_json: Any,
    template: Optional[str] = None,
    num_slides: Optional[int] = None,
    output_format: str = "pptx",
    ttl_minutes: int = 30,
) -> Dict[str, str]:
    if isinstance(slides_json, str):
        slides_obj = coerce_json(slides_json)
    else:
        slides_obj = slides_json or {}

    slides = slides_obj.get("slides", [])
    if num_slides:
        slides = slides[:num_slides]

    prs = Presentation(template) if template and os.path.exists(template) else Presentation()

    for s in slides:
        title: str = s.get("title", "")
        content: List[str] = s.get("content", []) or []
        notes: str = s.get("speaker_notes", "")
        image_path: Optional[str] = s.get("_image_path") or s.get("image_path") or None

        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title
        _add_bullets_to_slide(slide, content)
        _add_optional_image(slide, image_path)

        try:
            ns = slide.notes_slide
            ns.notes_text_frame.text = notes or ""
        except Exception:
            pass

    file_id = str(uuid.uuid4())
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    base = f"presentation_{timestamp}_{file_id}"

    output_format = output_format.lower()
    if output_format not in ["pptx", "pdf"]:
        raise ValueError("Unsupported format. Only 'pptx' and 'pdf' allowed.")

    if output_format == "pptx":
        pptx_path = os.path.join(OUT_DIR, base + ".pptx")
        prs.save(pptx_path)
        return save_and_record_export(pptx_path, output_format="pptx")

    pdf_path = os.path.join(OUT_DIR, base + ".pdf")
    width, height = landscape(A4)
    c = canvas.Canvas(pdf_path, pagesize=landscape(A4))

    for s in slides:
        title = s.get("title", "")
        content = s.get("content", []) or []

        c.setFont("Helvetica-Bold", 24)
        c.drawString(40, height - 80, title)
        c.setFont("Helvetica", 14)
        y = height - 120
        for p in content:
            c.drawString(60, y, u"\u2022 " + p)
            y -= 24
            if y < 60:
                c.showPage()
                c.setFont("Helvetica", 14)
                y = height - 120
        c.showPage()
    c.save()
    return save_and_record_export(pdf_path, output_format="pdf")